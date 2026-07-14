"""
Open GeoGuessr — Slide builder (korak po korak)
Pokretati svaki put kad se dodaje novi slajd.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

OUT = "/Users/lazarstoiljkovic/Desktop/other/open-geoguessr/OpenGeoGuessr_v2.pptx"

W = Inches(13.33)
H = Inches(7.5)

# ── Paleta (identična webapp scss/_variables.scss) ─────────────────────────
BG         = RGBColor(0x09, 0x0b, 0x14)   # $color-bg
SURFACE    = RGBColor(0x0f, 0x13, 0x20)   # $color-surface
SURFACE2   = RGBColor(0x14, 0x19, 0x29)   # $color-surface-2
SURFACE3   = RGBColor(0x1a, 0x20, 0x35)   # $color-surface-3
GOLD       = RGBColor(0xf0, 0xc0, 0x40)   # $color-primary
GOLD_DARK  = RGBColor(0xc8, 0x9f, 0x28)   # $color-primary-dark
BLUE       = RGBColor(0x5a, 0xab, 0xff)   # $color-accent
GREEN      = RGBColor(0x4c, 0xdd, 0x8a)   # $color-success
RED        = RGBColor(0xff, 0x5f, 0x5f)   # $color-danger
TEXT       = RGBColor(0xee, 0xf0, 0xf8)   # $color-text
MUTED      = RGBColor(0x7a, 0x85, 0xb0)   # $color-text-muted

# ── Helpers ─────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def load_or_create():
    if os.path.exists(OUT):
        prs = Presentation(OUT)
        prs.slide_width  = W
        prs.slide_height = H
    else:
        prs = new_prs()
    return prs

def add_slide(prs):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s

def shape_rect(slide, x, y, w, h, fill=SURFACE2, line=None, line_w=Pt(1)):
    from pptx.util import Pt
    sp = slide.shapes.add_shape(1, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line:
        sp.line.color.rgb = line
        sp.line.width = line_w
    else:
        sp.line.fill.background()
    return sp

def shape_line(slide, x1, y, x2, color=GOLD, width=Pt(1.5)):
    """Horizontalna linija."""
    sp = slide.shapes.add_shape(1, x1, y, x2 - x1, Pt(1.5))
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    sp.line.fill.background()
    return sp

def txt(slide, text, x, y, w, h,
        size=18, bold=False, color=TEXT,
        align=PP_ALIGN.LEFT, italic=False, font="Space Grotesk"):
    box = slide.shapes.add_textbox(x, y, w, h)
    box.word_wrap = True
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = font
    return box

def multiline_txt(slide, lines, x, y, w, h, default_size=14,
                  default_color=TEXT, default_font="DM Sans"):
    """lines = list of (text, size, bold, color, align, font) tuples."""
    box = slide.shapes.add_textbox(x, y, w, h)
    box.word_wrap = True
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        text  = item[0]
        size  = item[1] if len(item) > 1 else default_size
        bold  = item[2] if len(item) > 2 else False
        color = item[3] if len(item) > 3 else default_color
        align = item[4] if len(item) > 4 else PP_ALIGN.LEFT
        font  = item[5] if len(item) > 5 else default_font
        space = item[6] if len(item) > 6 else 0

        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        if space:
            p.space_before = Pt(space)
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = font
    return box

# ═══════════════════════════════════════════════════════════════════════════
# SLAJD 1 — NASLOVNI
# ═══════════════════════════════════════════════════════════════════════════
def slide_01(prs):
    s = add_slide(prs)

    # ── Decorative dot grid (gornjih 60% slajda) ───────────────────────────
    dot_color = RGBColor(0x1e, 0x26, 0x40)
    cols, rows = 42, 14
    dx = W / cols
    dy = (H * 0.62) / rows
    for r in range(rows):
        for c in range(cols):
            dot = s.shapes.add_shape(1,
                int(c * dx), int(r * dy),
                int(Inches(0.06)), int(Inches(0.06)))
            dot.fill.solid()
            dot.fill.fore_color.rgb = dot_color
            dot.line.fill.background()

    # ── Zlatni gradient overlay na dnu dot grida (simulacija) ─────────────
    fade = shape_rect(s, 0, Inches(4.2), W, Inches(1.1), fill=BG)
    fade.line.fill.background()

    # ── Zlatna tanka linija na vrhu ────────────────────────────────────────
    shape_line(s, Inches(0.6), Inches(0.15), W - Inches(0.6), GOLD, Pt(1))

    # ── Pin ikonica (jednostavna simulacija sa krugom + trouglom) ─────────
    pin_x = Inches(6.2)
    pin_y = Inches(1.35)
    circle = s.shapes.add_shape(9,  # oval
        pin_x, pin_y, Inches(0.9), Inches(0.9))
    circle.fill.solid()
    circle.fill.fore_color.rgb = GOLD
    circle.line.fill.background()

    inner = s.shapes.add_shape(9,
        pin_x + Inches(0.27), pin_y + Inches(0.27),
        Inches(0.36), Inches(0.36))
    inner.fill.solid()
    inner.fill.fore_color.rgb = BG
    inner.line.fill.background()

    # ── Naziv igre ─────────────────────────────────────────────────────────
    txt(s, "Open GeoGuessr",
        Inches(0.6), Inches(2.45), Inches(12.13), Inches(1.4),
        size=62, bold=True, color=TEXT,
        align=PP_ALIGN.CENTER, font="Space Grotesk")

    # ── Zlatna reč naglašena ────────────────────────────────────────────────
    # Overlay samo "Geo" delom — radimo kao poseban textbox koji prekriva
    # (simulacija jer pptx nema inline color u istom para bez run-a)
    # Koristimo drugi pristup: dve reči u kartici ispod naslova

    # ── Tagline ────────────────────────────────────────────────────────────
    txt(s, "Multiplayer geografska igra u realnom vremenu",
        Inches(0.6), Inches(3.75), Inches(12.13), Inches(0.55),
        size=16, bold=False, color=MUTED,
        align=PP_ALIGN.CENTER, font="DM Sans")

    # ── Zlatna separator linija ────────────────────────────────────────────
    shape_line(s, Inches(4.5), Inches(4.5), Inches(8.83), GOLD, Pt(1.5))

    # ── Info kartica ───────────────────────────────────────────────────────
    card = shape_rect(s,
        Inches(3.5), Inches(4.75),
        Inches(6.33), Inches(2.35),
        fill=SURFACE2,
        line=RGBColor(0xf0, 0xc0, 0x40),
        line_w=Pt(0.75))

    # Sadržaj kartice
    info_lines = [
        ("Lazar Stoiljković",   20, True,  TEXT,  PP_ALIGN.CENTER, "Space Grotesk", 0),
        ("Broj indeksa: 16381", 13, False, MUTED, PP_ALIGN.CENTER, "DM Sans", 10),
        ("",                    8,  False, MUTED, PP_ALIGN.CENTER, "DM Sans", 0),
        ("Arhitektura i projektovanje softvera",
                               12, False, BLUE,  PP_ALIGN.CENTER, "DM Sans", 6),
    ]
    multiline_txt(s, info_lines,
        Inches(3.6), Inches(4.92),
        Inches(6.13), Inches(2.0))

    # ── Zlatna linija na dnu ───────────────────────────────────────────────
    shape_line(s, Inches(0.6), H - Inches(0.2), W - Inches(0.6), GOLD, Pt(1))

    return s


# ═══════════════════════════════════════════════════════════════════════════
# SLAJD 2 — ŠTA JE IGRA + TOK IGRE
# ═══════════════════════════════════════════════════════════════════════════
def slide_02(prs):
    s = add_slide(prs)

    # Zlatna linija gore
    shape_line(s, Inches(0.6), Inches(0.12), W - Inches(0.6), GOLD, Pt(1))

    # Section badge
    badge = shape_rect(s, Inches(0.6), Inches(0.22), Inches(3.2), Inches(0.38),
                       fill=SURFACE3, line=GOLD, line_w=Pt(0.75))
    tf = badge.text_frame
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "PREGLED APLIKACIJE"
    r.font.size = Pt(10); r.font.bold = True; r.font.color.rgb = GOLD; r.font.name = "Space Grotesk"

    # Naslov
    txt(s, "Šta je Open GeoGuessr?",
        Inches(0.6), Inches(0.72), Inches(12.13), Inches(0.75),
        size=36, bold=True, color=TEXT, align=PP_ALIGN.LEFT, font="Space Grotesk")

    # Opis
    desc = shape_rect(s, Inches(0.6), Inches(1.55), Inches(12.13), Inches(0.72),
                      fill=SURFACE2, line=RGBColor(0x1e, 0x26, 0x40))
    tf = desc.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = ("Open GeoGuessr je open-source multiplayer geografska igra u realnom vremenu. "
              "Igrači gledaju fotografiju nepoznate lokacije putem Google Street View "
              "i pokušavaju da pogode tačno mesto na interaktivnoj mapi. "
              "Što bliže lokaciji — to više poena.")
    r.font.size = Pt(13); r.font.color.rgb = MUTED; r.font.name = "DM Sans"

    # ── Tok igre — label ───────────────────────────────────────────────────
    txt(s, "Tok igre",
        Inches(0.6), Inches(2.42), Inches(4.0), Inches(0.38),
        size=14, bold=True, color=GOLD, font="Space Grotesk")

    steps = [
        ("🏠", "Lobby",       "Host kreira sobu,\nigrači se pridružuju\nRoom Code-om"),
        ("⏳", "Countdown",   "5-sekundi odbrojavanje\npre svake runde"),
        ("🌍", "Street View", "Igrači razgledaju\nlokaciju i postavljaju\npin na mapu"),
        ("📍", "Guess",       "Interaktivna mapa,\nklik za postavljanje\npina"),
        ("📊", "Rezultati",   "Mapa sa gessovima\nsvih igrača,\nskorovi po rundi"),
        ("🏆", "Game Over",   "Pobednik, porijum,\ndetaljni skor\npo rundama"),
    ]

    # Izračunaj dimenzije da sve stane: W - 2*margin = 12.13"
    # 6 karata + 5 strelica: 6*sw + 5*aw = 12.13
    aw     = Inches(0.18)
    sw     = (Inches(12.13) - 5 * aw) / 6   # ≈ 1.873"
    bh     = Inches(1.62)
    bx0    = Inches(0.6)
    by     = Inches(2.88)

    # Fiksni y-offseti unutar kartice (apsolutni, isti za sve)
    OFF_NUM   = Inches(0.10)   # broj badge
    OFF_ICON  = Inches(0.10)   # ikona (pored broja)
    OFF_TITLE = Inches(0.54)   # naslov
    OFF_DESC  = Inches(0.88)   # opis

    for i, (icon, title, desc_text) in enumerate(steps):
        x = bx0 + i * (sw + aw)
        card_fill = SURFACE3 if i % 2 == 0 else SURFACE2
        border_col = GOLD if i == 0 else RGBColor(0x2a, 0x33, 0x50)
        border_w   = Pt(1.5) if i == 0 else Pt(0.75)

        # Kartica — čist pravougaonik, bez teksta
        shape_rect(s, x, by, sw, bh, fill=card_fill,
                   line=border_col, line_w=border_w)

        # Broj badge
        nb = shape_rect(s, x + Inches(0.1), by + OFF_NUM,
                        Inches(0.3), Inches(0.3),
                        fill=GOLD if i == 0 else SURFACE3, line=None)
        tf = nb.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = str(i + 1)
        r.font.size = Pt(10); r.font.bold = True
        r.font.color.rgb = BG if i == 0 else MUTED
        r.font.name = "Space Grotesk"

        # Ikona — isti red kao broj, odmah desno
        txt(s, icon,
            x + Inches(0.45), by + OFF_ICON,
            sw - Inches(0.5), Inches(0.38),
            size=18, color=TEXT, align=PP_ALIGN.LEFT)

        # Naslov — fiksan y za sve kartice
        txt(s, title,
            x + Inches(0.1), by + OFF_TITLE,
            sw - Inches(0.18), Inches(0.3),
            size=11, bold=True,
            color=GOLD if i == 0 else TEXT,
            align=PP_ALIGN.LEFT, font="Space Grotesk")

        # Opis — fiksan y za sve kartice
        txt(s, desc_text,
            x + Inches(0.1), by + OFF_DESC,
            sw - Inches(0.18), Inches(0.68),
            size=9, color=MUTED,
            align=PP_ALIGN.LEFT, font="DM Sans")

        # Strelica između karata
        if i < len(steps) - 1:
            ax = x + sw + Inches(0.01)
            ay = by + bh / 2 - Inches(0.2)
            txt(s, "›", ax, ay, aw, Inches(0.38),
                size=20, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    # ── Modovi igre — label ────────────────────────────────────────────────
    txt(s, "Modovi igre",
        Inches(0.6), Inches(4.67), Inches(4.0), Inches(0.38),
        size=14, bold=True, color=GOLD, font="Space Grotesk")

    modes = [
        ("🎮  Standard",   BLUE,  "Svi igraju N rundi.\nPobjeđuje najveći\nukupni skor."),
        ("⚡  Elimination", RED,   "Najgori igrač ispada\nposle svake runde.\nPobjeđuje poslednji."),
        ("👥  Teams",       GREEN, "2v2 timovi. Tim chat,\ntim skor, grupni\nresultati."),
        ("👁️  Spectator",   MUTED, "Gledaoci prate\nigru u realnom\nvremenu."),
    ]

    # 4 mode kartice: jednaka širina, fiksni margini
    mode_gap = Inches(0.18)
    mode_w   = (Inches(12.13) - 3 * mode_gap) / 4   # ≈ 2.97"
    mode_h   = Inches(1.93)
    mode_y   = Inches(5.1)

    # Fiksni y-offseti unutar mode kartice
    M_TITLE = Inches(0.15)
    M_DESC  = Inches(0.52)

    for i, (name, col, mdesc) in enumerate(modes):
        mx = Inches(0.6) + i * (mode_w + mode_gap)

        # Kartica
        shape_rect(s, mx, mode_y, mode_w, mode_h,
                   fill=SURFACE2, line=col, line_w=Pt(0.75))
        # Leva obojena ivica
        shape_rect(s, mx, mode_y, Inches(0.055), mode_h,
                   fill=col, line=None)

        # Naziv — fiksan y
        txt(s, name,
            mx + Inches(0.12), mode_y + M_TITLE,
            mode_w - Inches(0.15), Inches(0.32),
            size=12, bold=True, color=col,
            align=PP_ALIGN.LEFT, font="Space Grotesk")

        # Opis — fiksan y
        txt(s, mdesc,
            mx + Inches(0.12), mode_y + M_DESC,
            mode_w - Inches(0.15), Inches(1.3),
            size=10, color=MUTED,
            align=PP_ALIGN.LEFT, font="DM Sans")

    # Zlatna linija dole
    shape_line(s, Inches(0.6), H - Inches(0.15), W - Inches(0.6), GOLD, Pt(1))
    return s


# ── Shared helper za sve screen slajdove ────────────────────────────────────
def _screen_header(s, subtitle):
    shape_line(s, Inches(0.6), Inches(0.12), W - Inches(0.6), GOLD, Pt(1))
    badge = shape_rect(s, Inches(0.6), Inches(0.22), Inches(3.2), Inches(0.38),
                       fill=SURFACE3, line=GOLD, line_w=Pt(0.75))
    tf = badge.text_frame
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "EKRANI APLIKACIJE"
    r.font.size = Pt(10); r.font.bold = True; r.font.color.rgb = GOLD; r.font.name = "Space Grotesk"
    txt(s, subtitle,
        Inches(0.6), Inches(0.72), Inches(12.13), Inches(0.65),
        size=34, bold=True, color=TEXT, align=PP_ALIGN.LEFT, font="Space Grotesk")
    shape_line(s, Inches(0.6), H - Inches(0.15), W - Inches(0.6), GOLD, Pt(1))

def _ph(s, x, y, w, h):
    """Screenshot placeholder okvir."""
    frame = shape_rect(s, x, y, w, h, fill=SURFACE3,
                       line=RGBColor(0x2a, 0x33, 0x50), line_w=Pt(1.2))
    txt(s, "[ screenshot ]", x, y + h/2 - Inches(0.18), w, Inches(0.35),
        size=11, color=RGBColor(0x2a, 0x33, 0x50),
        align=PP_ALIGN.CENTER, font="Space Mono")
    return frame

def _desc_card(s, x, y, w, h, items):
    """Kartica sa bullet opisima desno od screenshota."""
    card = shape_rect(s, x, y, w, h, fill=SURFACE2,
                      line=RGBColor(0x2a, 0x33, 0x50), line_w=Pt(0.75))
    tf = card.text_frame; tf.word_wrap = True
    first = True
    for (label, desc, col) in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(0) if first else Pt(10)
        r = p.add_run(); r.text = label
        r.font.size = Pt(13); r.font.bold = True
        r.font.color.rgb = col; r.font.name = "Space Grotesk"
        p2 = tf.add_paragraph(); p2.space_before = Pt(4)
        r2 = p2.add_run(); r2.text = desc
        r2.font.size = Pt(11); r2.font.color.rgb = MUTED; r2.font.name = "DM Sans"


# ═══════════════════════════════════════════════════════════════════════════
# SLAJD 3a — LOBBY
# ═══════════════════════════════════════════════════════════════════════════
def slide_03a(prs):
    s = add_slide(prs)
    _screen_header(s, "Lobby — Čekaonica")

    # Screenshot levo
    _ph(s, Inches(0.6), Inches(1.5), Inches(8.2), Inches(5.6))

    # ── Desni panel — 4 tab-kartice ────────────────────────────────────────
    px      = Inches(9.05)
    py      = Inches(1.5)
    pw      = Inches(3.68)
    ph_tot  = Inches(5.6)
    gap     = Inches(0.12)
    n       = 4
    row_h   = (ph_tot - (n - 1) * gap) / n   # ≈ 1.31"

    tabs = [
        ("🔑", "Room Code",    GOLD,
         "Host kreira sobu i dobija unikalani 6-slovni Room Code "
         "(npr. SNMWSB). Igrači unose kod i odmah ulaze u čekaonicu."),
        ("👥", "Igrači",       BLUE,
         "Lista svih igrača u sobi se ažurira u realnom vremenu. "
         "Host nosi oznaku HOST. Igrač koji se diskonektuje prikazuje se zasivljeno."),
        ("⚙️", "Podešavanja",  GREEN,
         "Host bira game mode, broj rundi (1–10), trajanje runde i tip lokacija "
         "(Famous / Random World). Svi igrači vide podešavanja pre starta."),
        ("▶️", "Start Game",   GOLD,
         "Samo host može pokrenuti igru. Nakon klika kreće 5-sekundni countdown "
         "koji vide svi igrači, a zatim počinje prva runda."),
    ]

    for i, (icon, label, col, desc_text) in enumerate(tabs):
        ry = py + i * (row_h + gap)

        # Pozadina kartice
        shape_rect(s, px, ry, pw, row_h, fill=SURFACE2,
                   line=RGBColor(0x2a, 0x33, 0x50), line_w=Pt(0.75))

        # Leva obojena ivica
        shape_rect(s, px, ry, Inches(0.055), row_h, fill=col, line=None)

        # Icon + Label — u istom redu, top
        txt(s, icon,
            px + Inches(0.13), ry + Inches(0.13),
            Inches(0.38), Inches(0.38),
            size=16, color=TEXT, align=PP_ALIGN.LEFT)

        txt(s, label,
            px + Inches(0.55), ry + Inches(0.14),
            pw - Inches(0.65), Inches(0.36),
            size=15, bold=True, color=col,
            align=PP_ALIGN.LEFT, font="Space Grotesk")

        # Separator linija ispod naslova
        shape_line(s,
            px + Inches(0.13), ry + Inches(0.56),
            px + pw - Inches(0.1), col, Pt(0.5))

        # Opis — pun width kartice, ispod linije
        txt(s, desc_text,
            px + Inches(0.13), ry + Inches(0.64),
            pw - Inches(0.22), row_h - Inches(0.72),
            size=10.5, color=MUTED,
            align=PP_ALIGN.LEFT, font="DM Sans")

    return s


# ═══════════════════════════════════════════════════════════════════════════
# SLAJD 3b — STREET VIEW + PLACE PIN
# ═══════════════════════════════════════════════════════════════════════════
def _tab_caption(s, x, y, w, h, icon, label, col, desc_text):
    """Tab-kartica ispod screenshota — isti stil kao slide_03a desni panel."""
    shape_rect(s, x, y, w, h, fill=SURFACE2,
               line=RGBColor(0x2a, 0x33, 0x50), line_w=Pt(0.75))
    # Leva obojena ivica
    shape_rect(s, x, y, Inches(0.055), h, fill=col, line=None)
    # Ikona
    txt(s, icon,
        x + Inches(0.13), y + Inches(0.13),
        Inches(0.38), Inches(0.38),
        size=16, color=TEXT, align=PP_ALIGN.LEFT)
    # Label
    txt(s, label,
        x + Inches(0.55), y + Inches(0.14),
        w - Inches(0.65), Inches(0.36),
        size=15, bold=True, color=col,
        align=PP_ALIGN.LEFT, font="Space Grotesk")
    # Separator
    shape_line(s, x + Inches(0.13), y + Inches(0.57), x + w - Inches(0.1), col, Pt(0.5))
    # Opis
    txt(s, desc_text,
        x + Inches(0.13), y + Inches(0.65),
        w - Inches(0.22), h - Inches(0.73),
        size=10.5, color=MUTED,
        align=PP_ALIGN.LEFT, font="DM Sans")


def slide_03b(prs):
    s = add_slide(prs)
    _screen_header(s, "Gameplay")

    ph_w = Inches(5.8)
    ph_h = Inches(3.95)
    cap_h = Inches(1.45)
    gap   = Inches(0.3)
    y     = Inches(1.5)

    # Levi screenshot
    _ph(s, Inches(0.6), y, ph_w, ph_h)

    # Leva tab-kartica
    _tab_caption(s,
        Inches(0.6), y + ph_h + Inches(0.12),
        ph_w, cap_h,
        "🌍", "Razgledanje lokacije", GOLD,
        "Igrač dobija 360° fotografiju nepoznate lokacije putem Google Street View "
        "Može da se kreće i rotira. Opcioni hintovi otkrivaju kontinent ili državu lokacije.")

    # Desni screenshot
    rx = Inches(0.6) + ph_w + gap
    _ph(s, rx, y, ph_w, ph_h)

    # Desna tab-kartica
    _tab_caption(s,
        rx, y + ph_h + Inches(0.12),
        ph_w, cap_h,
        "📍", "Postavljanje pretpostavke", BLUE,
        "Igrač klikne na interaktivnu Google Maps mapu da postavi pin na pretpostavljenu "
        "lokaciju. Potvrdom gessa server računa distancu i dodeljuje poene.")

    return s


# ═══════════════════════════════════════════════════════════════════════════
# SLAJD 3c — REZULTATI RUNDE + GAME OVER
# ═══════════════════════════════════════════════════════════════════════════
def slide_03c(prs):
    s = add_slide(prs)
    _screen_header(s, "Rezultati & Game Over")

    ph_w  = Inches(5.8)
    ph_h  = Inches(3.95)
    cap_h = Inches(1.45)
    gap   = Inches(0.3)
    y     = Inches(1.5)

    # Levi screenshot
    _ph(s, Inches(0.6), y, ph_w, ph_h)
    _tab_caption(s,
        Inches(0.6), y + ph_h + Inches(0.12),
        ph_w, cap_h,
        "📊", "Rezultati runde", GOLD,
        "Po isteku runde prikazuje se mapa sa pinovima svih igrača i linijama do "
        "tačne lokacije. Svaki igrač vidi distancu i poene koje je zaradio u toj rundi.")

    # Desni screenshot
    rx = Inches(0.6) + ph_w + gap
    _ph(s, rx, y, ph_w, ph_h)
    _tab_caption(s,
        rx, y + ph_h + Inches(0.12),
        ph_w, cap_h,
        "🏆", "Game Over", GREEN,
        "Na kraju igre prikazuje se podijum sa top 3 igrača i detaljan skor po "
        "rundama za svakog učesnika. Host može pokrenuti novu igru klikom na Back to Lobby.")

    return s


# ═══════════════════════════════════════════════════════════════════════════
# SLAJD 3d — LEADERBOARD
# ═══════════════════════════════════════════════════════════════════════════
def slide_03d(prs):
    s = add_slide(prs)
    _screen_header(s, "Globalni Leaderboard")

    # Screenshot levo
    _ph(s, Inches(0.6), Inches(1.5), Inches(8.2), Inches(5.6))

    # Desni panel — 4 tab-kartice (isti stil kao slide_03a)
    px     = Inches(9.05)
    py     = Inches(1.5)
    pw     = Inches(3.68)
    ph_tot = Inches(5.6)
    gap    = Inches(0.12)
    n      = 4
    row_h  = (ph_tot - (n - 1) * gap) / n

    tabs = [
        ("🌍", "Globalni rang",  GOLD,
         "Svi registrovani igrači rangirani po ukupnom skoru "
         "sakupljenom kroz sve odigrane partije."),
        ("📋", "Kolone",         BLUE,
         "Rank · Player · Total Score · Games Played · Avg Score. "
         "Sortirano silazno po ukupnom skoru."),
        ("⚡", "Ažuriranje",     GREEN,
         "Skor se ažurira odmah po završetku svake igre putem "
         "Observer patterna (LeaderboardObserver)."),
        ("🥇", "Medalje",        GOLD,
         "Top 3 igrača dobijaju zlatnu, srebrnu i bronzanu medalju "
         "prikazane uz njihovo ime na listi."),
    ]

    for i, (icon, label, col, desc_text) in enumerate(tabs):
        ry = py + i * (row_h + gap)

        shape_rect(s, px, ry, pw, row_h, fill=SURFACE2,
                   line=RGBColor(0x2a, 0x33, 0x50), line_w=Pt(0.75))
        shape_rect(s, px, ry, Inches(0.055), row_h, fill=col, line=None)

        txt(s, icon,
            px + Inches(0.13), ry + Inches(0.13),
            Inches(0.38), Inches(0.38),
            size=16, color=TEXT, align=PP_ALIGN.LEFT)

        txt(s, label,
            px + Inches(0.55), ry + Inches(0.14),
            pw - Inches(0.65), Inches(0.36),
            size=15, bold=True, color=col,
            align=PP_ALIGN.LEFT, font="Space Grotesk")

        shape_line(s,
            px + Inches(0.13), ry + Inches(0.56),
            px + pw - Inches(0.1), col, Pt(0.5))

        txt(s, desc_text,
            px + Inches(0.13), ry + Inches(0.64),
            pw - Inches(0.22), row_h - Inches(0.72),
            size=10.5, color=MUTED,
            align=PP_ALIGN.LEFT, font="DM Sans")

    return s


# ═══════════════════════════════════════════════════════════════════════════
# SLAJD 4 — KOMPLETNA LISTA FEATURE-A
# ═══════════════════════════════════════════════════════════════════════════
def slide_04(prs):
    s = add_slide(prs)

    shape_line(s, Inches(0.6), Inches(0.12), W - Inches(0.6), GOLD, Pt(1))

    badge = shape_rect(s, Inches(0.6), Inches(0.22), Inches(3.2), Inches(0.38),
                       fill=SURFACE3, line=GOLD, line_w=Pt(0.75))
    tf = badge.text_frame
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "PREGLED APLIKACIJE"
    r.font.size = Pt(10); r.font.bold = True; r.font.color.rgb = GOLD; r.font.name = "Space Grotesk"

    txt(s, "Funkcionalnosti",
        Inches(0.6), Inches(0.72), Inches(12.13), Inches(0.6),
        size=36, bold=True, color=TEXT, align=PP_ALIGN.LEFT, font="Space Grotesk")

    cols_data = [
        {
            "title": "🎮  Gameplay",
            "color": GOLD,
            "items": [
                ("Multiplayer sobe",        "Host kreira sobu, igrači se pridružuju Room Code-om"),
                ("Standard mod",            "N rundi (1–10), pobjeđuje najveći ukupni skor"),
                ("Elimination mod",         "Najgori igrač ispada posle svake runde"),
                ("Famous Locations",        "39 poznatih svetskih lokacija"),
                ("World Locations",         "Nasumične lokacije putem Google Street View API-ja"),
                ("Countdown između rundi",  "3-sekundno odbrojavanje pre svake runde"),
                ("Podešavanje trajanja",    "Host bira trajanje runde: 30s – 5 min"),
            ],
        },
        {
            "title": "💬  Komunikacija",
            "color": BLUE,
            "items": [
                ("Globalni chat",        "Real-time chat za sve igrače u sobi"),
                ("Tim chat",             "Privatni chat samo za članove istog tima"),
                ("Live pin",             "Međusobno kretanje pina vidljivo u realnom vremenu"),
                ("Spectator mod",        "Gledaoci prate igru sa live pinovima svih igrača"),
                ("Room Code deljenje",   "Jednoklično kopiranje koda sobe"),
                ("Reconnect",           "Diskonektovani igrač može da se vrati u igru"),
            ],
        },
        {
            "title": "⚡  Napredne opcije",
            "color": GREEN,
            "items": [
                ("Timovi (2v2)",          "Dva tima po 2 igrača, grupni skor i pobednik"),
                ("Hints — Kontinent",     "Hint koji otkriva kontinent lokacije"),
                ("Hints — Država",        "Hint koji otkriva državu lokacije"),
                ("Hint kazna",            "Svaki hint smanjuje skor runde za 15%"),
                ("Bonus za preciznost",   "Guess unutar 100 km nosi +20% bonus"),
                ("Vremenski bonus",       "U Elimination modu: brži guess = više poena"),
                ("Globalni Leaderboard",  "Rang lista svih korisnika sa ukupnim skorom"),
            ],
        },
    ]

    col_w     = Inches(4.0)
    col_gap   = Inches(0.22)
    col_x0    = Inches(0.6)
    col_y     = Inches(1.55)
    col_h     = Inches(5.7)
    pad_l     = Inches(0.13)   # padding od leve ivice (posle accent bara)
    pad_r     = Inches(0.12)
    accent_w  = Inches(0.055)

    # Fiksni y-offseti unutar kartice
    OFF_ICON  = Inches(0.13)
    OFF_TITLE = Inches(0.14)
    OFF_SEP   = Inches(0.56)
    OFF_ITEMS = Inches(0.65)

    for ci, col in enumerate(cols_data):
        cx = col_x0 + ci * (col_w + col_gap)

        # Pozadina kartice — čist pravougaonik
        shape_rect(s, cx, col_y, col_w, col_h,
                   fill=SURFACE2, line=col["color"], line_w=Pt(0.75))

        # Leva obojena ivica
        shape_rect(s, cx, col_y, accent_w, col_h, fill=col["color"], line=None)

        # Ikona naslova
        txt(s, col["title"].split("  ")[0],
            cx + pad_l, col_y + OFF_ICON,
            Inches(0.38), Inches(0.36),
            size=16, color=TEXT, align=PP_ALIGN.LEFT)

        # Naslov
        txt(s, col["title"].split("  ")[1],
            cx + pad_l + Inches(0.42), col_y + OFF_TITLE,
            col_w - pad_l - Inches(0.42) - pad_r, Inches(0.36),
            size=15, bold=True, color=col["color"],
            align=PP_ALIGN.LEFT, font="Space Grotesk")

        # Separator linija
        shape_line(s,
            cx + pad_l, col_y + OFF_SEP,
            cx + col_w - pad_r, col["color"], Pt(0.5))

        # Stavke — jedan textbox koji počinje od vrha (textbox, ne shape)
        items_box = s.shapes.add_textbox(
            cx + pad_l, col_y + OFF_ITEMS,
            col_w - pad_l - pad_r,
            col_h - OFF_ITEMS - Inches(0.1))
        items_box.word_wrap = True
        tf = items_box.text_frame
        tf.word_wrap = True

        first = True
        for feat_name, feat_desc in col["items"]:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            if not first:
                p.space_before = Pt(8)
            r = p.add_run()
            r.text = "▸  " + feat_name
            r.font.size = Pt(13)
            r.font.bold = True
            r.font.color.rgb = TEXT
            r.font.name = "DM Sans"

            p2 = tf.add_paragraph()
            p2.space_before = Pt(2)
            r2 = p2.add_run()
            r2.text = "    " + feat_desc
            r2.font.size = Pt(11)
            r2.font.color.rgb = MUTED
            r2.font.name = "DM Sans"

    shape_line(s, Inches(0.6), H - Inches(0.15), W - Inches(0.6), GOLD, Pt(1))
    return s


# ═══════════════════════════════════════════════════════════════════════════
# Shared: tech kartice
# ═══════════════════════════════════════════════════════════════════════════
def _fe_badge(s, x, y, w, h, icon, name, col):
    """Badge samo sa ikonom i nazivom, bez opisa — vertikalno centriran sadržaj."""
    shape_rect(s, x, y, w, h, fill=SURFACE2,
               line=RGBColor(0x22, 0x2d, 0x45), line_w=Pt(0.75))
    shape_rect(s, x, y, Inches(0.055), h, fill=col, line=None)
    txt(s, icon,
        x + Inches(0.18), y + h / 2 - Inches(0.22),
        Inches(0.44), Inches(0.44),
        size=22, color=TEXT, align=PP_ALIGN.LEFT)
    txt(s, name,
        x + Inches(0.7), y + h / 2 - Inches(0.2),
        w - Inches(0.8), Inches(0.42),
        size=18, bold=True, color=col,
        align=PP_ALIGN.LEFT, font="Space Grotesk")


def _tech_card(s, x, y, w, h, icon, name, col, desc):
    shape_rect(s, x, y, w, h, fill=SURFACE2,
               line=RGBColor(0x22, 0x2d, 0x45), line_w=Pt(0.75))
    shape_rect(s, x, y, Inches(0.055), h, fill=col, line=None)

    txt(s, icon,
        x + Inches(0.13), y + Inches(0.14),
        Inches(0.38), Inches(0.38),
        size=18, color=TEXT, align=PP_ALIGN.LEFT)

    txt(s, name,
        x + Inches(0.58), y + Inches(0.15),
        w - Inches(0.68), Inches(0.38),
        size=15, bold=True, color=col,
        align=PP_ALIGN.LEFT, font="Space Grotesk")

    shape_line(s, x + Inches(0.13), y + Inches(0.6),
               x + w - Inches(0.12), col, Pt(0.5))

    txt(s, desc,
        x + Inches(0.13), y + Inches(0.68),
        w - Inches(0.22), h - Inches(0.76),
        size=11, color=MUTED,
        align=PP_ALIGN.LEFT, font="DM Sans")


def _stack_header(s, badge_text, title, title_color=TEXT):
    shape_line(s, Inches(0.6), Inches(0.12), W - Inches(0.6), GOLD, Pt(1))
    badge = shape_rect(s, Inches(0.6), Inches(0.22), Inches(3.2), Inches(0.38),
                       fill=SURFACE3, line=GOLD, line_w=Pt(0.75))
    tf = badge.text_frame
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = badge_text
    r.font.size = Pt(10); r.font.bold = True; r.font.color.rgb = GOLD; r.font.name = "Space Grotesk"
    txt(s, title,
        Inches(0.6), Inches(0.72), Inches(12.13), Inches(0.65),
        size=36, bold=True, color=title_color,
        align=PP_ALIGN.LEFT, font="Space Grotesk")
    shape_line(s, Inches(0.6), H - Inches(0.15), W - Inches(0.6), GOLD, Pt(1))


# ═══════════════════════════════════════════════════════════════════════════
# SLAJD 5a — BACKEND
# ═══════════════════════════════════════════════════════════════════════════
def slide_05a(prs):
    s = add_slide(prs)
    _stack_header(s, "TEHNOLOGIJE", "Backend")

    items = [
        ("🟢", "Node.js + Koa",       GOLD,
         "HTTP server i REST API.\n"
         "Rute za autentifikaciju, kreiranje soba, tok igre i leaderboard.\n"
         "Koa middleware lanac za validaciju i autentifikaciju zahteva."),
        ("🔷", "TypeScript",          BLUE,
         "Statičko tipiziranje kroz ceo backend — interfejsi, generici, union tipovi.\n"
         "Kompajler hvata greške pre pokretanja: pogrešni tipovi, nedostajući argumenti."),
        ("💉", "TypeDI",              GREEN,
         "Dependency Injection kontejner. @Service dekorator označava klasu kao singleton.\n"
         "Framework automatski injektuje zavisnosti u konstruktor — nema ručnog new()."),
        ("🔌", "WebSocket (ws)",      GOLD,
         "Native WebSocket server, bez Socket.IO. Konekcija se čuva dok god je klijent aktivan.\n"
         "JWT autentifikacija na samoj konekciji (token u URL query stringu)."),
        ("🔑", "jsonwebtoken",        BLUE,
         "JWT tokeni za autentifikaciju HTTP ruta i WebSocket konekcija.\n"
         "Token se verifikuje pri svakom zahtevu i pri WS handshake-u."),
        ("🌐", "axios",               GREEN,
         "HTTP klijent za pozive ka eksternim servisima:\n"
         "Google Street View Static API (verifikacija panorama) i BigDataCloud (reverse geocoding za hintove)."),
    ]

    card_w = Inches(5.85)
    card_h = Inches(1.48)
    gap_x  = Inches(0.28)
    gap_y  = Inches(0.16)
    x0     = Inches(0.6)
    y0     = Inches(1.55)

    for i, (icon, name, col, desc) in enumerate(items):
        col_i = i % 2
        row_i = i // 2
        x = x0 + col_i * (card_w + gap_x)
        y = y0 + row_i * (card_h + gap_y)
        _tech_card(s, x, y, card_w, card_h, icon, name, col, desc)

    return s


# ═══════════════════════════════════════════════════════════════════════════
# SLAJD 5b — BAZA + INFRASTRUKTURA
# ═══════════════════════════════════════════════════════════════════════════
def slide_05b(prs):
    s = add_slide(prs)
    _stack_header(s, "TEHNOLOGIJE", "Baza podataka & Infrastruktura")

    # ── Levo: Baza podataka ───────────────────────────────────────────────
    bx = Inches(0.6)
    bw = Inches(5.85)
    by = Inches(1.55)

    txt(s, "🗄️  Baza podataka", bx, by, bw, Inches(0.4),
        size=16, bold=True, color=BLUE, font="Space Grotesk")
    shape_line(s, bx, by + Inches(0.45), bx + bw, BLUE, Pt(0.75))

    db_items = [
        ("🍃", "MongoDB",       BLUE,
         "NoSQL document baza podataka.\n"
         "3 kolekcije: Users (korisnici), Rooms (sobe sa rundama i gessovima), "
         "CachedLocations (potvrđene Street View lokacije).\n"
         "Embedded arrays za runde, gessove i chat poruke — bez JOIN operacija."),
        ("📐", "Mongoose ODM",  BLUE,
         "Object-Document Mapper za TypeScript.\n"
         "Svaka kolekcija ima Schema (tipovi, default vrednosti, validacija) i "
         "TypeScript interface (IRoom, IUser, ICachedLocation).\n"
         "Domenski metodi enkapsulisani u Repository klase — servisi ne vide Mongoose direktno."),
    ]

    # leva strana: 2 karte sa opisima, puna visina
    card_h = Inches(2.45)
    gap    = Inches(0.18)
    cy = by + Inches(0.55)
    for icon, name, col, desc in db_items:
        _tech_card(s, bx, cy, bw, card_h, icon, name, col, desc)
        cy += card_h + gap

    # ── Desno: Infrastruktura ─────────────────────────────────────────────
    ix = Inches(6.68)
    iw = Inches(6.0)
    iy = Inches(1.55)

    txt(s, "🚀  Infrastruktura", ix, iy, iw, Inches(0.4),
        size=16, bold=True, color=GREEN, font="Space Grotesk")
    shape_line(s, ix, iy + Inches(0.45), ix + iw, GREEN, Pt(0.75))

    # desna strana: samo ikona + naziv (bez opisa) — koristi _fe_badge
    infra_items = [
        ("🖥️", "Hetzner CX22",   GREEN),
        ("🐳", "Docker Compose", GREEN),
        ("🔁", "Nginx",          GREEN),
        ("🔒", "Let's Encrypt",  GREEN),
        ("▲",  "Vercel",         BLUE),
    ]

    # 5 kartica unutar iste visine kao leva strana (5.08")
    card_h2 = (Inches(5.08) - 4 * Inches(0.12)) / 5
    gap2    = Inches(0.12)
    cy2 = iy + Inches(0.55)
    for icon, name, col in infra_items:
        _fe_badge(s, ix, cy2, iw, card_h2, icon, name, col)
        cy2 += card_h2 + gap2

    return s


# ═══════════════════════════════════════════════════════════════════════════
# SLAJD 5c — FRONTEND
# ═══════════════════════════════════════════════════════════════════════════
def slide_05c(prs):
    s = add_slide(prs)
    _stack_header(s, "TEHNOLOGIJE", "Frontend")

    fe_items = [
        ("⚛️",  "React 18",                RED),
        ("🔷", "TypeScript",               BLUE),
        ("🔄", "React Context + useReducer", RED),
        ("🎨", "SCSS Modules",             BLUE),
        ("🗺️", "Google Maps JS API",       GREEN),
        ("🔌", "WebSocket (browser native)", GOLD),
    ]

    # 3 kolone × 2 reda
    card_w = Inches(3.86)
    card_h = Inches(2.42)
    gap_x  = Inches(0.28)
    gap_y  = Inches(0.25)
    x0     = Inches(0.6)
    y0     = Inches(1.55)

    for i, (icon, name, col) in enumerate(fe_items):
        col_i = i % 3
        row_i = i // 3
        x = x0 + col_i * (card_w + gap_x)
        y = y0 + row_i * (card_h + gap_y)
        _fe_badge(s, x, y, card_w, card_h, icon, name, col)

    return s


# ── Code block + Collection card helpers ───────────────────────────────────
def _code_block(s, x, y, w, h, lines, size=9):
    shape_rect(s, x, y, w, h, fill=RGBColor(0x07, 0x09, 0x10),
               line=RGBColor(0x2a, 0x33, 0x50), line_w=Pt(1))
    box = s.shapes.add_textbox(x + Inches(0.15), y + Inches(0.15),
                               w - Inches(0.3), h - Inches(0.3))
    box.word_wrap = False
    tf = box.text_frame; tf.word_wrap = False
    first = True
    for text, col in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        r = p.add_run(); r.text = text
        r.font.size = Pt(size); r.font.name = "Space Mono"; r.font.color.rgb = col
    return box


def _uml_box(s, x, y, w, h, stereotype, name, fields, col):
    """UML klasa/dokument: obojena header sekcija + lista polja."""
    hdr_h = Inches(0.55)
    shape_rect(s, x, y, w, h, fill=SURFACE2, line=col, line_w=Pt(1))
    shape_rect(s, x, y, w, hdr_h, fill=SURFACE3, line=None)
    shape_rect(s, x, y, Inches(0.045), h, fill=col, line=None)
    txt(s, f"«{stereotype}»",
        x + Inches(0.06), y + Inches(0.04), w - Inches(0.06), Inches(0.18),
        size=8, color=col, align=PP_ALIGN.CENTER, italic=True, font="Space Grotesk")
    txt(s, name,
        x + Inches(0.06), y + Inches(0.23), w - Inches(0.06), Inches(0.28),
        size=12, bold=True, color=col, align=PP_ALIGN.CENTER, font="Space Grotesk")
    shape_line(s, x, y + hdr_h, x + w, col, Pt(0.75))
    fbox = s.shapes.add_textbox(
        x + Inches(0.1), y + hdr_h + Inches(0.07),
        w - Inches(0.15), h - hdr_h - Inches(0.1))
    fbox.word_wrap = False
    tf = fbox.text_frame; tf.word_wrap = False
    first = True
    for line in fields:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        r = p.add_run(); r.text = line
        r.font.size = Pt(8.5); r.font.name = "Space Mono"
        r.font.color.rgb = MUTED if (line.startswith('  ') or line.startswith('─')) else TEXT


def _h_arrow(s, x1, y, x2, label, col=MUTED):
    """Horizontalna strelica x1→x2 sa labelom iznad."""
    shape_line(s, x1, y, x2 - Inches(0.1), col, Pt(1.2))
    txt(s, "▶", x2 - Inches(0.14), y - Inches(0.13),
        Inches(0.16), Inches(0.22), size=8.5, color=col, align=PP_ALIGN.LEFT)
    if label:
        txt(s, label, x1 + Inches(0.02), y - Inches(0.23),
            x2 - x1 - Inches(0.1), Inches(0.2),
            size=7, color=col, italic=True, align=PP_ALIGN.CENTER, font="Space Mono")


def _coll_card(s, x, y, w, h, icon, name, col, fields):
    shape_rect(s, x, y, w, h, fill=SURFACE2,
               line=RGBColor(0x22, 0x2d, 0x45), line_w=Pt(0.75))
    shape_rect(s, x, y, Inches(0.055), h, fill=col, line=None)
    txt(s, icon, x + Inches(0.12), y + Inches(0.12),
        Inches(0.35), Inches(0.35), size=13, color=TEXT)
    txt(s, name, x + Inches(0.52), y + Inches(0.14),
        w - Inches(0.6), Inches(0.32),
        size=12, bold=True, color=col, align=PP_ALIGN.LEFT, font="Space Grotesk")
    shape_line(s, x + Inches(0.12), y + Inches(0.52), x + w - Inches(0.1), col, Pt(0.5))
    bx = s.shapes.add_textbox(x + Inches(0.12), y + Inches(0.60), w - Inches(0.2), h - Inches(0.68))
    bx.word_wrap = False
    tf = bx.text_frame; tf.word_wrap = False
    first = True
    for line in fields:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        is_special = line.startswith(' ') or line.startswith('─')
        if ': ' in line and not is_special:
            fn, ft = line.split(': ', 1)
            r1 = p.add_run(); r1.text = fn
            r1.font.size = Pt(8.5); r1.font.bold = True
            r1.font.name = "Space Mono"; r1.font.color.rgb = TEXT
            r2 = p.add_run(); r2.text = f': {ft}'
            r2.font.size = Pt(8.5); r2.font.name = "Space Mono"; r2.font.color.rgb = MUTED
        else:
            r = p.add_run(); r.text = line
            r.font.size = Pt(8.5); r.font.name = "Space Mono"; r.font.color.rgb = MUTED


# ═══════════════════════════════════════════════════════════════════════════
# SLAJD 06a — ZAŠTO MONGODB?
# ═══════════════════════════════════════════════════════════════════════════
def slide_06a(prs):
    s = add_slide(prs)
    _stack_header(s, "PERZISTENCIJA PODATAKA", "Zašto MongoDB?")

    y0, th, gap = Inches(1.55), Inches(5.65), Inches(0.16)
    lx, lw = Inches(0.6),  Inches(5.93)
    rx, rw = Inches(6.78), Inches(5.95)

    lh = (th - 2 * gap) / 3
    for i, (icon, label, col, desc) in enumerate([
        ("📄", "Embedded Documents", GOLD,
         "Cela igra živi u jednom Room dokumentu: igrači, runde, gessovi i chat.\n"
         "Čitanje stanja igre = jedan MongoDB upit, bez JOIN-ova."),
        ("⚡", "Atomski array operatori", BLUE,
         "$push, $set + arrayFilters, $addToSet rade direktno na ugn. nizovima.\n"
         "Guess se upisuje u rounds[i].guesses[] jednim atomskim updateom."),
        ("🔧", "Fleksibilna šema", GREEN,
         "Nema SQL migracija. Novo polje (teamsEnabled) → dodamo u Schema → deploy.\n"
         "SQL: ALTER TABLE + migraciona skripta za svaki strukturni change."),
    ]):
        _tab_caption(s, lx, y0 + i * (lh + gap), lw, lh, icon, label, col, desc)

    rh = (th - gap) / 2
    for i, (icon, label, col, desc) in enumerate([
        ("⚠️", "Ograničenja", RED,
         "Bez ACID transakcija po defaultu — update više dokumenata nije atomski.\n"
         "Nema referencijalnog integriteta: userId u Round.guess može biti invalid.\n"
         "Kompleksne agregacije su neprirodnijе nego u SQL-u."),
        ("🤔", "Kad SQL ima smisla?", MUTED,
         "Finansijski sistemi, billing — gde su konzistentnost i JOIN upiti\n"
         "između entiteta kritični. Za real-time game use-case MongoDB je\n"
         "prirodniji izbor jer celokupno stanje igre je jedan dokument."),
    ]):
        _tab_caption(s, rx, y0 + i * (rh + gap), rw, rh, icon, label, col, desc)

    return s


# ═══════════════════════════════════════════════════════════════════════════
# SLAJD 06b — ŠEMA BAZE PODATAKA
# ═══════════════════════════════════════════════════════════════════════════
def slide_06b(prs):
    s = add_slide(prs)
    _stack_header(s, "PERZISTENCIJA PODATAKA", "Šema baze podataka")

    y0 = Inches(1.55)
    h  = Inches(5.65)

    x_left  = Inches(0.6);   w_left  = Inches(2.3)
    x_room  = Inches(3.2);   w_room  = Inches(3.3)   # gap 0.3"
    x_emb   = Inches(7.0);   w_emb   = Inches(2.3)   # gap 0.5" (za strelice)
    x_guess = Inches(9.65);  w_guess = Inches(2.48)  # gap 0.35"

    # ── Levo: User + CachedLocation ───────────────────────────────────────
    _uml_box(s, x_left, y0, w_left, Inches(2.55), "collection", "User", [
        "+ _id: ObjectId",
        "+ username: String {unique}",
        "+ email: String {unique}",
        "+ password: String",
        "+ totalScore: Number",
        "+ gamesPlayed: Number",
        "─" * 13,
        "+ createdAt: Date",
        "+ updatedAt: Date",
    ], BLUE)

    _uml_box(s, x_left, Inches(4.28), w_left, Inches(2.92), "collection", "CachedLocation", [
        "+ _id: ObjectId",
        "+ lat: Number",
        "+ lng: Number",
        "+ panoId: String",
        "+ name: String",
        "+ country: String",
        "+ usedCount: Number",
        "  {index: true}",
        "+ lastUsedAt?: Date",
        "─" * 13,
        "+ createdAt: Date",
    ], GREEN)

    # ── Centar: Room ──────────────────────────────────────────────────────
    _uml_box(s, x_room, y0, w_room, h, "collection", "Room", [
        "+ _id: ObjectId",
        "+ code: String {unique, UPPER}",
        "+ hostId: String",
        "+ status: GameStatus",
        "  {enum, 5 states}",
        "+ gameMode: GameMode",
        "  {standard | elimination}",
        "+ locationMode: LocationMode",
        "  {famous | world}",
        "+ totalRounds: Number",
        "+ roundDurationSeconds: Number",
        "+ currentRoundIndex: Number",
        "+ eliminatedPlayerIds: String[]",
        "+ hintsEnabled: Boolean",
        "+ spectatorsAllowed: Boolean",
        "+ teamsEnabled: Boolean",
        "+ teamSize: Number",
        "─" * 15,
        "  players[]:   →  Player",
        "  rounds[]:    →  Round",
        "  messages[]:  →  ChatMessage",
        "─" * 15,
        "+ createdAt: Date",
        "+ updatedAt: Date",
    ], GOLD)

    # ── Desno: Embedded dokumenti ─────────────────────────────────────────
    y_player = y0
    h_player = Inches(1.3)
    _uml_box(s, x_emb, y_player, w_emb, h_player, "embedded", "Player", [
        "+ userId: String",
        "+ username: String",
        "+ isHost: Boolean",
        "+ score: Number",
        "+ connected: Boolean",
        "+ teamId?: Number",
    ], BLUE)

    y_round = y_player + h_player + Inches(0.18)
    h_round = Inches(1.75)
    _uml_box(s, x_emb, y_round, w_emb, h_round, "embedded", "Round", [
        "+ index: Number",
        "+ location: Location{}",
        "+ startedAt: Number",
        "+ endedAt?: Number",
        "─" * 11,
        "  guesses[]:  →  RoundGuess",
    ], GREEN)

    y_msg = y_round + h_round + Inches(0.18)
    h_msg = (y0 + h) - y_msg
    _uml_box(s, x_emb, y_msg, w_emb, h_msg, "embedded", "ChatMessage", [
        "+ userId: String",
        "+ username: String",
        "+ text: String",
        "+ timestamp: Number",
    ], MUTED)

    # ── Krajnje desno: RoundGuess ─────────────────────────────────────────
    _uml_box(s, x_guess, y_round, w_guess, h_round, "embedded", "RoundGuess", [
        "+ userId: String",
        "+ lat: Number",
        "+ lng: Number",
        "+ distanceKm: Number",
        "+ roundScore: Number",
        "+ submittedAt: Number",
    ], GREEN)

    # ── Composition strelice ──────────────────────────────────────────────
    room_r = x_room + w_room
    _h_arrow(s, room_r, y_player + h_player / 2, x_emb, "players[]  1..*", BLUE)
    _h_arrow(s, room_r, y_round  + h_round  / 2, x_emb, "rounds[]   1..*", GREEN)
    _h_arrow(s, room_r, y_msg    + h_msg    / 2, x_emb, "messages[] 0..*", MUTED)
    _h_arrow(s, x_emb + w_emb, y_round + h_round / 2, x_guess, "guesses[]  0..*", GREEN)

    return s


# ═══════════════════════════════════════════════════════════════════════════
# SLAJD 06c — MONGOOSE ODM ARHITEKTURA
# ═══════════════════════════════════════════════════════════════════════════
def slide_06c(prs):
    s = add_slide(prs)
    _stack_header(s, "PERZISTENCIJA PODATAKA", "Mongoose ODM")

    y0 = Inches(1.55)
    h  = Inches(5.65)

    w1 = Inches(3.4);  x1 = Inches(0.6)
    w2 = Inches(3.4);  x2 = Inches(4.7)    # gap 0.7"
    w3 = Inches(3.33); x3 = Inches(8.8)    # gap 0.7"
    # x3 + w3 = 12.13 ✓

    # Box 1 — TypeScript Interface
    _uml_box(s, x1, y0, w1, h, "interface", "IRoom", [
        "+ code: string",
        "+ hostId: string",
        "+ status: GameStatus",
        "+ gameMode: GameMode",
        "+ locationMode: LocationMode",
        "+ players: Player[]",
        "+ rounds: Round[]",
        "+ messages: ChatMessage[]",
        "+ eliminatedPlayerIds: string[]",
        "+ totalRounds: number",
        "+ roundDurationSeconds: number",
        "+ currentRoundIndex: number",
        "+ hintsEnabled: boolean",
        "+ spectatorsAllowed: boolean",
        "+ teamsEnabled: boolean",
        "+ teamSize: number",
        "─" * 15,
        "  extends Document",
    ], BLUE)

    # Box 2 — Mongoose Schema
    _uml_box(s, x2, y0, w2, h, "schema", "RoomSchema", [
        "code: String",
        "  {required, unique, UPPER}",
        "status: String",
        "  {enum: 5 states}",
        "  {default: 'waiting'}",
        "gameMode: String",
        "  {enum: standard|elimination}",
        "locationMode: String",
        "  {enum: famous|world}",
        "players:  [PlayerSchema]",
        "rounds:   [RoundSchema]",
        "messages: [ChatMsgSchema]",
        "eliminatedPlayerIds: [String]",
        "  {default: []}",
        "hintsEnabled: Boolean",
        "  {default: false}",
        "─" * 16,
        "  {timestamps: true}",
    ], GOLD)

    # Box 3 — MongoDB Collection
    _uml_box(s, x3, y0, w3, h, "document", "rooms", [
        "_id: ObjectId  (auto)",
        "code: \"SNMWSB\"",
        "status: \"playing\"",
        "gameMode: \"standard\"",
        "locationMode: \"famous\"",
        "totalRounds: 5",
        "currentRoundIndex: 2",
        "hintsEnabled: false",
        "spectatorsAllowed: false",
        "teamsEnabled: false",
        "─" * 14,
        "players: Array(3)",
        "rounds: Array(5)",
        "messages: Array(12)",
        "eliminatedPlayerIds: []",
        "─" * 14,
        "createdAt: ISODate",
        "updatedAt: ISODate",
    ], GREEN)

    # Horizontalne strelice između boxova
    arr_y = y0 + h / 2
    _h_arrow(s, x1 + w1, arr_y, x2, "Schema<IRoom>",    BLUE)
    _h_arrow(s, x2 + w2, arr_y, x3, "mongoose.model()", GOLD)

    return s


# ═══════════════════════════════════════════════════════════════════════════
# SLAJD 06d — REPOSITORY PATTERN (UML class dijagram)
# ═══════════════════════════════════════════════════════════════════════════
def slide_06d(prs):
    s = add_slide(prs)
    _stack_header(s, "PERZISTENCIJA PODATAKA", "Repository Pattern")

    y0   = Inches(1.55)
    h    = Inches(5.65)
    bx   = Inches(0.6)
    bw   = Inches(12.13)
    hdr  = Inches(0.55)
    cgap = Inches(0.25)
    cw   = (bw - Inches(0.4) - cgap) / 2   # ≈ 5.74"

    h_base = Inches(2.45)
    h_gap  = Inches(0.3)
    h_room = h - h_base - h_gap             # ≈ 2.9"

    # ── Helper: 2-kolumnski textbox unutar UML boxa ──────────────────────
    def _col(cx, cy, cw_, ch, lines):
        box = s.shapes.add_textbox(cx, cy, cw_, ch)
        box.word_wrap = False; tf = box.text_frame; tf.word_wrap = False
        first = True
        for line in lines:
            p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
            r = p.add_run(); r.text = line
            r.font.size = Pt(9); r.font.name = "Space Mono"
            r.font.color.rgb = MUTED if (
                line.startswith('─') or line.startswith('  ')
                or line.startswith('#') or line == ""
            ) else TEXT

    def _uml_wide(bx_, by_, bh_, col_, stereo, name_, left_fields, right_fields):
        shape_rect(s, bx_, by_, bw, bh_, fill=SURFACE2, line=col_, line_w=Pt(1))
        shape_rect(s, bx_, by_, bw, hdr,  fill=SURFACE3, line=None)
        shape_rect(s, bx_, by_, Inches(0.045), bh_, fill=col_, line=None)
        txt(s, f"«{stereo}»",
            bx_ + Inches(0.06), by_ + Inches(0.04), bw - Inches(0.06), Inches(0.18),
            size=8, color=col_, align=PP_ALIGN.CENTER, italic=True, font="Space Grotesk")
        txt(s, name_,
            bx_ + Inches(0.06), by_ + Inches(0.23), bw - Inches(0.06), Inches(0.28),
            size=13, bold=True, color=col_, align=PP_ALIGN.CENTER, font="Space Grotesk")
        shape_line(s, bx_, by_ + hdr, bx_ + bw, col_, Pt(0.75))
        fy = by_ + hdr + Inches(0.08)
        fh = bh_ - hdr - Inches(0.12)
        _col(bx_ + Inches(0.15),          fy, cw, fh, left_fields)
        _col(bx_ + Inches(0.15) + cw + cgap, fy, cw, fh, right_fields)

    # ── BaseRepository<T> ─────────────────────────────────────────────────
    _uml_wide(bx, y0, h_base, BLUE, "abstract", "BaseRepository<T extends Document>",
        left_fields=[
            "# model: Model<T>  (protected readonly)",
            "─" * 34,
            "+ findById(id: string): Promise<T | null>",
            "+ findOne(filter: FilterQuery<T>): Promise<T | null>",
            "+ findMany(filter: FilterQuery<T>): Promise<T[]>",
            "+ findPaginated(filter, opts): Promise<PaginatedResult<T>>",
            "+ create(data: Partial<T>): Promise<T>",
        ],
        right_fields=[
            "",
            "─" * 34,
            "+ updateById(id, data: UpdateQuery<T>): Promise<T | null>",
            "+ updateOne(filter, data: UpdateQuery<T>): void",
            "+ deleteById(id: string): void",
            "+ deleteMany(filter: FilterQuery<T>): void",
            "+ exists(filter: FilterQuery<T>): boolean",
            "+ count(filter: FilterQuery<T>): number",
        ]
    )

    # ── Inheritance strelica (vertikalna, gore) ───────────────────────────
    arr_x = bx + bw / 2
    shape_rect(s, arr_x - Inches(0.015), y0 + h_base, Inches(0.03), h_gap, fill=BLUE, line=None)
    txt(s, "▲", arr_x - Inches(0.1), y0 + h_base - Inches(0.06),
        Inches(0.2), Inches(0.2), size=9, color=BLUE, align=PP_ALIGN.LEFT)
    txt(s, "extends", arr_x + Inches(0.05), y0 + h_base + Inches(0.05),
        Inches(0.7), Inches(0.18), size=7.5, color=BLUE, italic=True,
        align=PP_ALIGN.LEFT, font="Space Grotesk")

    # ── RoomRepository ────────────────────────────────────────────────────
    _uml_wide(bx, y0 + h_base + h_gap, h_room, GOLD,
        "@Service() · repository", "RoomRepository",
        left_fields=[
            "  extends BaseRepository<IRoom>",
            "─" * 34,
            "+ findByCode(code: string): Promise<IRoom | null>",
            "+ createRoom(data: {...}): Promise<IRoom>",
            "+ addPlayer(roomId, player): Promise<IRoom | null>",
            "+ updatePlayerConnection(roomId, userId, bool): void",
            "+ setStatus(roomId, status): Promise<IRoom | null>",
            "+ initGame(roomId, rounds, ...): Promise<IRoom | null>",
            "+ setRoundLocation(roomId, idx, location): void",
        ],
        right_fields=[
            "",
            "─" * 34,
            "+ updateCurrentRound(roomId, idx, ...): Promise<IRoom | null>",
            "+ addGuessToRound(roomId, idx, guess): Promise<IRoom | null>",
            "+ updatePlayerScore(roomId, userId, score): void",
            "+ setRoundEndTime(roomId, idx, endedAt): void",
            "+ addEliminatedPlayer(roomId, userId): void",
            "+ updatePlayerTeam(roomId, userId, teamId): void",
            "+ addMessage(roomId, msg): void",
            "+ deleteOldRooms(): void",
        ]
    )


# ═══════════════════════════════════════════════════════════════════════════
# Shared helperi — arhitektura slojeva i projektni obrasci
# ═══════════════════════════════════════════════════════════════════════════
def _layer_band(s, x, y, w, h, title, col, subtitle=""):
    shape_rect(s, x, y, w, h, fill=SURFACE2, line=col, line_w=Pt(1))
    shape_rect(s, x, y, Inches(0.06), h, fill=col, line=None)
    txt(s, title, x + Inches(0.22), y + Inches(0.08), w - Inches(0.4), Inches(0.3),
        size=13, bold=True, color=col, font="Space Grotesk")
    if subtitle:
        txt(s, subtitle, x + Inches(0.22), y + Inches(0.4), w - Inches(0.4), h - Inches(0.5),
            size=9.5, color=MUTED, font="DM Sans")


def _v_chevron(s, x_center, y, col=MUTED):
    txt(s, "▼", x_center - Inches(0.15), y - Inches(0.02), Inches(0.3), Inches(0.16),
        size=9, bold=True, color=col, align=PP_ALIGN.CENTER)


def _flow_step(s, x, y, w, h, num, text, col=GOLD):
    shape_rect(s, x, y, w, h, fill=SURFACE2, line=RGBColor(0x2a, 0x33, 0x50), line_w=Pt(0.75))
    nb = shape_rect(s, x + Inches(0.12), y + h / 2 - Inches(0.17), Inches(0.34), Inches(0.34),
                    fill=col, line=None)
    tf = nb.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(num)
    r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = BG; r.font.name = "Space Grotesk"
    txt(s, text, x + Inches(0.6), y + h / 2 - Inches(0.26), w - Inches(0.75), h - Inches(0.1),
        size=10.5, color=TEXT, font="DM Sans")


def _table_row(s, x, y, cols, fill, header=False):
    """cols = list of (text, width, color)."""
    w = sum(cw for _, cw, _ in cols)
    h = Inches(0.5) if header else Inches(0.72)
    shape_rect(s, x, y, w, h, fill=fill, line=RGBColor(0x22, 0x2d, 0x45), line_w=Pt(0.5))
    cx = x
    for text, cw, col in cols:
        txt(s, text, cx + Inches(0.15), y + h / 2 - Inches(0.16), cw - Inches(0.25), Inches(0.4),
            size=12 if header else 11, bold=header, color=col,
            align=PP_ALIGN.LEFT, font="Space Grotesk" if header else "DM Sans")
        cx += cw
    return h


# ═══════════════════════════════════════════════════════════════════════════
# SLAJD 07a — VIŠESLOJNA ARHITEKTURA (BACKEND)
# ═══════════════════════════════════════════════════════════════════════════
def slide_07a(prs):
    s = add_slide(prs)
    _stack_header(s, "ARHITEKTURA APLIKACIJE", "Backend — slojevita arhitektura")

    x, w = Inches(0.6), Inches(8.7)
    y0 = Inches(1.55)
    h = Inches(0.7)
    gap = Inches(0.14)

    layers = [
        ("Klijent (Browser)", MUTED,
         "React SPA — stranice, komponente, GameContext (WS + HTTP klijent)"),
        ("Transport sloj", GOLD,
         "HTTP REST (Koa)  +  WebSocket (ws) — dele isti Node.js port"),
        ("Routing / Handler sloj", BLUE,
         "auth · rooms · game · leaderboard rute  /  GameHandler rutira WS evente"),
        ("Middleware sloj", RED,
         "JWT autentifikacija, validacija zahteva, error handling"),
        ("Servisni sloj — biznis logika", GREEN,
         "AuthService · RoomService · GameService · LocationService — koriste obrasce"),
        ("Repository sloj — perzistencija", GOLD,
         "BaseRepository<T> · RoomRepository · UserRepository · CachedLocationRepository"),
        ("Baza podataka", BLUE,
         "MongoDB + Mongoose ODM — kolekcije: users, rooms, cachedlocations"),
    ]

    ys = []
    y = y0
    for title, col, sub in layers:
        _layer_band(s, x, y, w, h, title, col, sub)
        ys.append(y)
        y += h + gap

    for i in range(len(layers) - 1):
        _v_chevron(s, x + Inches(0.6), ys[i] + h + Inches(0.01), MUTED)

    total_h = (ys[-1] + h) - y0
    nx = x + w + Inches(0.25)
    nw = W - nx - Inches(0.6)
    shape_rect(s, nx, y0, nw, total_h, fill=SURFACE3, line=GOLD, line_w=Pt(0.75))
    txt(s, "💉  Dependency Injection", nx + Inches(0.15), y0 + Inches(0.15), nw - Inches(0.3), Inches(0.55),
        size=13, bold=True, color=GOLD, align=PP_ALIGN.LEFT, font="Space Grotesk")
    txt(s,
        "TypeDI @Service() kontejner injektuje zavisnosti kroz sve slojeve: "
        "Routes/Handler → Services → Repositories.\n\n"
        "Svaka klasa dobija svoje zavisnosti kroz konstruktor — nema ručnog new().\n\n"
        "Zato je svaki sloj testabilan izolovano — mock repozitorijum se ubaci "
        "umesto pravog bez izmene servisa.",
        nx + Inches(0.15), y0 + Inches(0.75), nw - Inches(0.3), total_h - Inches(0.9),
        size=10.5, color=MUTED, align=PP_ALIGN.LEFT, font="DM Sans")

    return s


# ═══════════════════════════════════════════════════════════════════════════
# SLAJD 07b — FRONTEND SLOJEVI + TOK PODATAKA KROZ APLIKACIJU
# ═══════════════════════════════════════════════════════════════════════════
def slide_07b(prs):
    s = add_slide(prs)
    _stack_header(s, "ARHITEKTURA APLIKACIJE", "Frontend slojevi & tok podataka")

    lx, lw = Inches(0.6), Inches(5.7)
    y0 = Inches(1.55)
    lh_total = Inches(5.65)
    n = 4
    gap = Inches(0.16)
    lh = (lh_total - (n - 1) * gap) / n

    fe_layers = [
        ("Stranice (Pages)", BLUE,
         "LobbyPage · GamePage · LeaderboardPage — React Router"),
        ("Komponente", GOLD,
         "StreetViewPanel · GuessMap · Chat · Scoreboard · TeamPanel"),
        ("State sloj — Context + useReducer", GREEN,
         "GameContext drži GameState; menja se isključivo kroz dispatch(action)"),
        ("Transport klijenti", RED,
         "SocketClient (WS singleton, auto-reconnect) · gameApi / authApi (HTTP)"),
    ]
    ly = y0
    lys = []
    for title, col, sub in fe_layers:
        _layer_band(s, lx, ly, lw, lh, title, col, sub)
        lys.append(ly)
        ly += lh + gap
    for i in range(len(fe_layers) - 1):
        _v_chevron(s, lx + Inches(0.6), lys[i] + lh + Inches(0.01), MUTED)

    rx, rw = Inches(6.6), Inches(6.13)
    txt(s, "Tok podataka — primer: slanje pogotka (submitGuess)",
        rx, y0, rw, Inches(0.4), size=13, bold=True, color=GOLD, font="Space Grotesk")

    steps = [
        "GuessMap (komponenta) → gameApi.submitGuess(lat, lng)  — HTTP POST",
        "game.routes → GameService.submitGuess()",
        "ScoringContext.calculate() — Strategy + Decorator obrasci",
        "RoomRepository.addGuessToRound() — atomski $push u MongoDB",
        "HTTP odgovor → GameContext: dispatch({ type: 'GUESS_RESULT' })",
        "Kad svi pogode → broadcastToRoom('round_ended') — WebSocket svima",
    ]
    sy = y0 + Inches(0.5)
    sh_total = lh_total - Inches(0.5)
    sgap = Inches(0.1)
    sh = (sh_total - (len(steps) - 1) * sgap) / len(steps)
    for i, step in enumerate(steps):
        _flow_step(s, rx, sy + i * (sh + sgap), rw, sh, i + 1, step, GOLD)

    return s


def _problem_diagram_slide(prs, title, problem_text, diagram_lines):
    s = add_slide(prs)
    _stack_header(s, "PROJEKTNI OBRASCI", title)
    _tab_caption(s, Inches(0.6), Inches(1.55), Inches(12.13), Inches(1.3),
                 "❓", "Problem", RED, problem_text)
    _code_block(s, Inches(0.6), Inches(3.0), Inches(12.13), Inches(4.2), diagram_lines, size=12)
    return s


def _code_pros_cons_slide(prs, title, code_lines, pros_text, cons_text):
    s = add_slide(prs)
    _stack_header(s, "PROJEKTNI OBRASCI", title + " — implementacija")
    _code_block(s, Inches(0.6), Inches(1.55), Inches(12.13), Inches(4.2), code_lines, size=12)
    hw = (Inches(12.13) - Inches(0.3)) / 2
    _tab_caption(s, Inches(0.6), Inches(5.9), hw, Inches(1.3), "✅", "Prednosti", GREEN, pros_text)
    _tab_caption(s, Inches(0.6) + hw + Inches(0.3), Inches(5.9), hw, Inches(1.3),
                 "⚠️", "Mane", RED, cons_text)
    return s


# ═══════════════════════════════════════════════════════════════════════════
# SLAJD 08a — FACTORY METHOD (2 slajda)
# ═══════════════════════════════════════════════════════════════════════════
def slide_08a1(prs):
    return _problem_diagram_slide(prs, "🏭  Factory Method — izbor provajdera lokacija",
        "GameService treba lokaciju za svaku rundu iz dva izvora — Famous (statička lista) ili "
        "World (Google Street View API). Bez obrasca bi servis sam proveravao if/else i poznavao "
        "obe implementacije direktno.",
        [
            ("«interface» ILocationProvider", GOLD),
            ("  + getLocation(): Promise<Location>", TEXT),
            ("", MUTED),
            ("        ↑                        ↑", MUTED),
            ("FamousLocationProvider     WorldLocationProvider", BLUE),
            ("", MUTED),
            ("«abstract» LocationProviderFactory", GOLD),
            ("  + createProvider(mode): ILocationProvider", TEXT),
            ("      ← factory metoda (subklasa odlučuje)", MUTED),
            ("  + getLocation(mode): Promise<Location>", TEXT),
            ("      ← poziva createProvider() pa getLocation()", MUTED),
            ("        ↑", MUTED),
            ("GameLocationProviderFactory", GREEN),
            ("  ← konkretni kreator, zna koji provider da napravi", MUTED),
        ])


def slide_08a2(prs):
    return _code_pros_cons_slide(prs, "🏭  Factory Method",
        [
            ("// location-provider.factory.ts", MUTED),
            ("abstract class LocationProviderFactory {", TEXT),
            ("  abstract createProvider(mode): ILocationProvider;", BLUE),
            ("", MUTED),
            ("  async getLocation(mode) {", TEXT),
            ("    const provider = this.createProvider(mode);", GOLD),
            ("    return provider.getLocation();", TEXT),
            ("  }", TEXT),
            ("}", TEXT),
            ("", MUTED),
            ("@Service()", GREEN),
            ("class GameLocationProviderFactory extends LocationProviderFactory {", TEXT),
            ("  createProvider(mode) {", TEXT),
            ("    switch (mode) {", TEXT),
            ("      case 'famous': return this.famous;", MUTED),
            ("      case 'world':  return this.world;", MUTED),
            ("    }", TEXT),
            ("  }", TEXT),
            ("}", TEXT),
        ],
        "GameService ne zna ništa o FamousLocationProvider ni WorldLocationProvider. Dodavanje "
        "trećeg moda (npr. capitals) traži samo novu klasu + jedan case u createProvider().",
        "Dodaje slojeve indirekcije — debug prati lanac GameService → LocationService → Factory "
        "→ Provider. Ako postoji samo jedan provider, fabrika je over-engineering.")


# ═══════════════════════════════════════════════════════════════════════════
# SLAJD 08b — STRATEGY (2 slajda)
# ═══════════════════════════════════════════════════════════════════════════
def slide_08b1(prs):
    return _problem_diagram_slide(prs, "🧩  Strategy — algoritmi bodovanja",
        "Formula za bodovanje zavisi od game moda: Standard koristi samo distancu, Elimination "
        "dodaje vremenski bonus za brzinu. Bez obrasca GameService bi imao "
        "if (gameMode === 'elimination') {...} else {...} razbacano kroz kod.",
        [
            ("«interface» ScoringStrategy", GOLD),
            ("  + calculate(input: ScoringInput): number", TEXT),
            ("", MUTED),
            ("        ↑                          ↑", MUTED),
            ("DistanceScoringStrategy     TimeBonusScoringStrategy", BLUE),
            ("  (samo distanca)              (distanca + brzina)", MUTED),
            ("", MUTED),
            ("ScoringContext", GREEN),
            ("  - strategy: ScoringStrategy", TEXT),
            ("  + setStrategy(strategy)", TEXT),
            ("  + calculate(input): number", TEXT),
            ("      → this.strategy.calculate(input)", MUTED),
            ("", MUTED),
            ("Kreira se JEDNOM pri startGame() i čuva kroz celu partiju", MUTED),
            ("(in-memory mapa: Map<roomCode, ScoringContext>)", MUTED),
        ])


def slide_08b2(prs):
    return _code_pros_cons_slide(prs, "🧩  Strategy",
        [
            ("// distance.strategy.ts — Standard mod", MUTED),
            ("class DistanceScoringStrategy {", TEXT),
            ("  calculate({ distanceKm }) {", TEXT),
            ("    const ratio = Math.max(0, 1 - distanceKm / MAX_DISTANCE_KM);", TEXT),
            ("    return Math.round(MAX_SCORE * ratio ** 2);", GOLD),
            ("  }", TEXT),
            ("}  // 0km→5000  10000km→1250  20000km→0", MUTED),
            ("", MUTED),
            ("// time-bonus.strategy.ts — Elimination mod", MUTED),
            ("class TimeBonusScoringStrategy {", TEXT),
            ("  calculate({ distanceKm, timeTaken, roundDuration }) {", TEXT),
            ("    // 70% distanca + 30% brzina", MUTED),
            ("    return Math.round(MAX_SCORE * (", GOLD),
            ("      0.7 * distRatio ** 2 + 0.3 * timeRatio));", GOLD),
            ("  }", TEXT),
            ("}", TEXT),
        ],
        "Svaki mod ima nezavisnu, testabilnu formulu. ScoringContext ne zna koja se formula "
        "koristi — samo poziva calculate(). Nova formula = nova klasa, bez izmene GameService.",
        "Za samo 2 strategije korist je manja — isplati se tek kad broj varijanti raste. Klijent "
        "(GameService) i dalje mora da zna KOJU strategiju bira pri startGame().")


# ═══════════════════════════════════════════════════════════════════════════
# SLAJD 08c — DECORATOR (GAMESETUP) (2 slajda)
# ═══════════════════════════════════════════════════════════════════════════
def slide_08c1(prs):
    return _problem_diagram_slide(prs, "🎭  Decorator — konfiguracija sobe (GameSetup)",
        "Soba može biti kombinacija: Standard/Elimination × Hints × Teams × Spectators = 8 "
        "kombinacija opcija. Nasleđivanjem bi trebalo 8 klasa (npr. StandardGameWithHintsAndTeams) "
        "— svaki novi feature duplira broj klasa.",
        [
            ("«interface» IGameSetup", GOLD),
            ("  + describe(): string", TEXT),
            ("  + buildConfig(): GameSettings", TEXT),
            ("", MUTED),
            ("        ↑", MUTED),
            ("«abstract» GameSetupDecorator", BLUE),
            ("  # game: IGameSetup   (umotana instanca)", TEXT),
            ("", MUTED),
            ("        ↑              ↑              ↑", MUTED),
            ("HintsDecorator   TeamsDecorator   SpectatorDecorator", GREEN),
            ("", MUTED),
            ("StandardGameSetup  /  EliminationGameSetup", GOLD),
            ("  ← bazni, konkretni IGameSetup", MUTED),
        ])


def slide_08c2(prs):
    return _code_pros_cons_slide(prs, "🎭  Decorator (GameSetup)",
        [
            ("// room.service.ts — createRoom()", MUTED),
            ("let gameSetup: IGameSetup = gameMode === 'elimination'", TEXT),
            ("  ? new EliminationGameSetup()", GOLD),
            ("  : new StandardGameSetup();", GOLD),
            ("", MUTED),
            ("// slaganje dekoratora — uslovno, po jedan", MUTED),
            ("if (hintsEnabled)", TEXT),
            ("  gameSetup = new HintsDecorator(gameSetup);", GREEN),
            ("if (spectatorsAllowed)", TEXT),
            ("  gameSetup = new SpectatorDecorator(gameSetup);", GREEN),
            ("if (teamsEnabled)", TEXT),
            ("  gameSetup = new TeamsDecorator(gameSetup, teamSize);", GREEN),
            ("", MUTED),
            ("const settings = gameSetup.buildConfig();", GOLD),
            ("// describe() → 'Elimination + Hints + Teams (2v2)'", MUTED),
        ],
        "8 kombinacija bez 8 klasa — svaki dekorator dodaje samo svoju izmenu. Novi feature = "
        "jedna nova klasa, bez menjanja postojećih. describe() se lančano slaže za log/debug.",
        "Redosled umotavanja može biti bitan ako dekoratori nisu nezavisni. Debug zahteva praćenje "
        "celog lanca umotavanja da bi se video finalni config.")


# ═══════════════════════════════════════════════════════════════════════════
# SLAJD 08d — DECORATOR (SCORING) (2 slajda)
# ═══════════════════════════════════════════════════════════════════════════
def slide_08d1(prs):
    return _problem_diagram_slide(prs, "🎭  Decorator — modifikatori bodova (Scoring)",
        "Hint penal (-15% po hintu) i accuracy bonus (+20% na <100km) treba primeniti povrh bazne "
        "formule, za oba game moda. Nasleđivanjem bi trebalo DistanceWithHintPenaltyAndBonus, "
        "TimeBonusWithHintPenaltyAndBonus... kombinatorijalna eksplozija.",
        [
            ("«abstract» ScoringDecorator", GOLD),
            ("  # strategy: ScoringStrategy   (umotana)", TEXT),
            ("", MUTED),
            ("        ↑                          ↑", MUTED),
            ("HintPenaltyDecorator        AccuracyBonusDecorator", GREEN),
            ("  ×(1-0.15)^hintsUsed          ×1.20 ako distanca < 100km", MUTED),
            ("", MUTED),
            ("Primer: Elimination, 2 hinta, pogodak 50km", BLUE),
            ("", MUTED),
            ("TimeBonusScoringStrategy.calculate()   → 3200", TEXT),
            ("  HintPenaltyDecorator: ×0.72 (2 hinta) → 2304", GOLD),
            ("  AccuracyBonusDecorator: 50km < 100 → ×1.20", GOLD),
            ("  FINALNI SKOR → 2764", GREEN),
        ])


def slide_08d2(prs):
    return _code_pros_cons_slide(prs, "🎭  Decorator (Scoring)",
        [
            ("// hint-penalty.decorator.ts", MUTED),
            ("class HintPenaltyDecorator extends ScoringDecorator {", TEXT),
            ("  calculate(input) {", TEXT),
            ("    const base = this.strategy.calculate(input);", TEXT),
            ("    const mult = (1 - 0.15) ** (input.hintsUsed ?? 0);", GOLD),
            ("    return Math.round(base * mult);", GOLD),
            ("  }", TEXT),
            ("}", TEXT),
            ("", MUTED),
            ("// slaganje u GameService.startGame()", MUTED),
            ("let strategy = gameMode === 'elimination'", TEXT),
            ("  ? new TimeBonusScoringStrategy()", TEXT),
            ("  : new DistanceScoringStrategy();", TEXT),
            ("if (room.hintsEnabled)", TEXT),
            ("  strategy = new HintPenaltyDecorator(strategy);", GREEN),
            ("strategy = new AccuracyBonusDecorator(strategy);", GREEN),
        ],
        "Svaki dekorator ima jednu odgovornost i testira se izolovano. Kombinacije penala i "
        "bonusa menjaju se u runtime-u bez ijedne nove klase.",
        "Redosled slaganja je bitan — AccuracyBonusDecorator mora biti spolja jer se primenjuje "
        "na finalni skor. Greška u formuli zahteva praćenje celog lanca poziva.")


# ═══════════════════════════════════════════════════════════════════════════
# SLAJD 08e — OBSERVER (2 slajda)
# ═══════════════════════════════════════════════════════════════════════════
def slide_08e1(prs):
    return _problem_diagram_slide(prs, "👁️  Observer — obaveštenja pri kraju igre",
        "Kad igra završi treba: (1) emitovati game_over event svim igračima i (2) ažurirati "
        "globalni leaderboard u bazi. Bez obrasca GameService bi direktno zavisio od "
        "UserRepository-ja i WebSocket infrastrukture.",
        [
            ("«interface» GameObserver", GOLD),
            ("  + onGameOver(event): Promise<void>", TEXT),
            ("", MUTED),
            ("        ↑                          ↑", MUTED),
            ("LeaderboardObserver          BroadcastObserver", GREEN),
            ("  ažurira UserRepository        broadcastToRoom(WS)", MUTED),
            ("", MUTED),
            ("«abstract» GameEventSubject", BLUE),
            ("  - observers: GameObserver[]", TEXT),
            ("  + registerObserver(obs)", TEXT),
            ("  # notifyGameOver(event) → Promise.all(observers.map(...))", TEXT),
            ("        ↑", MUTED),
            ("GameService extends GameEventSubject", GOLD),
        ])


def slide_08e2(prs):
    return _code_pros_cons_slide(prs, "👁️  Observer",
        [
            ("// game.service.ts", MUTED),
            ("class GameService extends GameEventSubject {", TEXT),
            ("  constructor(leaderboardObs, broadcastObs) {", TEXT),
            ("    super();", TEXT),
            ("    this.registerObserver(leaderboardObs);", GREEN),
            ("    this.registerObserver(broadcastObs);", GREEN),
            ("  }", TEXT),
            ("  private async triggerGameOver(roomId, code) {", TEXT),
            ("    await this.notifyGameOver({ roomId, ... });", GOLD),
            ("  }", TEXT),
            ("}", TEXT),
            ("", MUTED),
            ("class LeaderboardObserver implements GameObserver {", TEXT),
            ("  async onGameOver({ players }) {", TEXT),
            ("    await Promise.all(players.map(p =>", TEXT),
            ("      userRepo.updateScore(p.userId, p.score)));", GOLD),
            ("  }", TEXT),
            ("}", TEXT),
        ],
        "GameService ne zna da postoje UserRepository ni WebSocket — samo za GameObserver "
        "interfejs. Novi feature (npr. email posle igre) = nova klasa + registerObserver().",
        "Promise.all odbacuje ceo poziv ako jedan observer baci grešku — u produkciji svaki "
        "observer treba da hvata sopstvene greške. Redosled notifikacija nije garantovan.")


# ═══════════════════════════════════════════════════════════════════════════
# SLAJD 08f — DEPENDENCY INJECTION (TypeDI) (2 slajda)
# ═══════════════════════════════════════════════════════════════════════════
def slide_08f1(prs):
    return _problem_diagram_slide(prs, "💉  Dependency Injection — TypeDI",
        "Svaki servis zavisi od više drugih klasa (repozitorijumi, drugi servisi, observeri). "
        "Ručno kreiranje (new GameService(new RoomRepository(), ...)) čvrsto veže klase i "
        "otežava testiranje — mock se ne može lako ubaciti.",
        [
            ("IoC kontejner (TypeDI)", GOLD),
            ("", MUTED),
            ("@Service()  →  markira klasu kao singleton", TEXT),
            ("", MUTED),
            ("Container.get(GameService)", GREEN),
            ("  → kreira/vraća singleton instancu", MUTED),
            ("  → automatski injektuje sve zavisnosti iz konstruktora, rekurzivno", MUTED),
            ("", MUTED),
            ("GameService", BLUE),
            ("  ← RoomRepository", MUTED),
            ("  ← LocationService ← GameLocationProviderFactory", MUTED),
            ("  ← LeaderboardObserver", MUTED),
            ("  ← BroadcastObserver", MUTED),
        ])


def slide_08f2(prs):
    return _code_pros_cons_slide(prs, "💉  Dependency Injection",
        [
            ("@Service()", GREEN),
            ("export class GameService extends GameEventSubject {", TEXT),
            ("  constructor(", TEXT),
            ("    private readonly roomRepository: RoomRepository,", BLUE),
            ("    private readonly locationService: LocationService,", BLUE),
            ("    leaderboardObserver: LeaderboardObserver,", BLUE),
            ("    broadcastObserver: BroadcastObserver,", BLUE),
            ("  ) { ... }", TEXT),
            ("}", TEXT),
            ("", MUTED),
            ("// u rutama:", MUTED),
            ("const gameService = () => Container.get(GameService);", GOLD),
            ("", MUTED),
            ("// BEZ DI (anti-pattern):", RED),
            ("const roomRepo = new RoomRepository();", MUTED),
            ("const gameService = new GameService(roomRepo, ...);", MUTED),
        ],
        "Nema ručnog new() — kontejner rešava ceo graf zavisnosti automatski. Testovi injektuju "
        "mock umesto pravog repozitorijuma bez izmene servisa.",
        "Nije klasičan GoF obrazac — magija dekoratora otežava praćenje odakle instanca dolazi "
        "ako se ne poznaje TypeDI. Sve zavisnosti moraju biti markirane @Service().")


# ═══════════════════════════════════════════════════════════════════════════
# SLAJD 08g — PREGLED SVIH OBRAZACA
# ═══════════════════════════════════════════════════════════════════════════
def slide_08g(prs):
    s = add_slide(prs)
    _stack_header(s, "PROJEKTNI OBRASCI", "Pregled svih obrazaca")

    x = Inches(0.6)
    c1, c2, c3 = Inches(2.75), Inches(2.55), Inches(6.83)

    rows = [
        ("Factory Method", "patterns/factory/",
         "Izbor FamousLocationProvider vs WorldLocationProvider po LocationMode", GOLD),
        ("Decorator (GameSetup)", "patterns/game-setup/",
         "Kombinovanje opcija sobe (Hints, Spectators, Teams) bez nasleđivanja", GREEN),
        ("Decorator (Scoring)", "patterns/scoring/",
         "Hint penal i accuracy bonus povrh bazne scoring strategije", GREEN),
        ("Strategy", "patterns/scoring/",
         "Zamenjivi algoritmi bodovanja (Distance, TimeBonus)", BLUE),
        ("Observer", "patterns/observer/",
         "LeaderboardObserver i BroadcastObserver reaguju na kraj igre", RED),
        ("Repository", "database/repositories/",
         "BaseRepository<T> + RoomRepository — apstrakcija nad Mongoose", GOLD),
        ("Dependency Injection", "ceo backend",
         "TypeDI @Service() — upravljanje lifecycle-om i zavisnostima servisa", BLUE),
    ]

    y = Inches(1.55)
    y += _table_row(s, x, y, [
        ("Obrazac", c1, GOLD), ("Lokacija", c2, GOLD), ("Uloga", c3, GOLD),
    ], fill=SURFACE3, header=True)

    for i, (name, loc, role, col) in enumerate(rows):
        fill = SURFACE2 if i % 2 == 0 else SURFACE
        y += _table_row(s, x, y, [
            (name, c1, col), (loc, c2, MUTED), (role, c3, TEXT),
        ], fill=fill)

    return s


# ── Main ────────────────────────────────────────────────────────────────────
prs = new_prs()
slide_01(prs)
slide_02(prs)
slide_03a(prs)
slide_03b(prs)
slide_03c(prs)
slide_03d(prs)
slide_04(prs)
slide_05a(prs)
slide_05b(prs)
slide_05c(prs)
slide_06a(prs)
slide_06b(prs)
slide_06c(prs)
slide_06d(prs)
slide_07a(prs)
slide_07b(prs)
slide_08a1(prs)
slide_08a2(prs)
slide_08b1(prs)
slide_08b2(prs)
slide_08c1(prs)
slide_08c2(prs)
slide_08d1(prs)
slide_08d2(prs)
slide_08e1(prs)
slide_08e2(prs)
slide_08f1(prs)
slide_08f2(prs)
slide_08g(prs)
prs.save(OUT)
print(f"✓  Saved ({len(prs.slides)} slide/s): {OUT}")
